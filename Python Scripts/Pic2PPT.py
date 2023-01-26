#pip install python-pptx==0.6.18
import os
import pyperclip
from pptx import Presentation
from pptx.util import Inches


# Get the folder path from the clipboard
folder_path = os.path.join(pyperclip.paste().strip('"'),'')

# Create a new PowerPoint file
prs = Presentation()

# Iterate through all the files in the folder
allowed_extensions = [".jpg", ".png"]
for file_name in os.listdir(folder_path):
    if file_name.endswith(tuple(allowed_extensions)):
        # Add the image to the PowerPoint file
        #replace the number 6 with any other index of the slide_layouts list to use a different layout for your slide.
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        #set the position of the image
        left = Inches(0)
        top = Inches(0)
        pic = slide.shapes.add_picture(os.path.join(folder_path, file_name), left, top)
        width  = 10
        height = 7.5
        pic.width = Inches(height)
        pic.height = Inches(height)







# Save the PowerPoint file in the same folder
prs.save(os.path.join(folder_path, "MergedImages.pptx"))
print("PowerPoint file saved successfully!")