# Written by Ahmad Cooper
# Generates flashcards in PowerPoint from a word list. You can copy the word list in CSV format or separated by new lines.
# Saves as PPT and exports to PDF

import collections.abc
import comtypes.client
from pptx import Presentation
import pyperclip
import re
import os
script_dir = os.path.dirname(os.path.abspath(__file__))
# If you don't change the directory, you can't run the script properly from a .bat file
os.chdir(script_dir)
# 
clipboard = pyperclip.paste()
words = re.split(',|\n', clipboard)
final = [word.strip() for word in words if word.strip()]

# Define the name of the PowerPoint file
FILENAME = "PyFC"
ext = ".pptx"
try:
    # Try to open an existing PowerPoint presentation or create a new one
    prs = Presentation(f"{FILENAME}{ext}")
    

    # Get the custom layout by ID
    custom_layout = prs.slide_layouts.get_by_name('3')
    if custom_layout is None:
        raise ValueError('Custom layout not found')
    

    for i in range(len(final)):
        # Add a new slide with the custom layout
        slide = prs.slides.add_slide(custom_layout)

        # Find the text box placeholder and set its text
        for shape in slide.placeholders:
            if shape.name == 'Text Placeholder 1':
                shape.text = str(final[i])
                break

        # Set the slide name
        slide.name = str(final[i])

    # Save the presentation
    prs.save(f"{FILENAME}_1{ext}")  

    powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
    powerpoint.Presentations.Open(os.path.abspath(f"{FILENAME}_1{ext}"))
    powerpoint.ActivePresentation.SaveAs(os.path.abspath(f"{FILENAME}_1.pdf"), 
                                      32) # 32 is the PDF format code
    powerpoint.Quit()
    print(f"Successfully created {len(final)} FC slides")

except FileNotFoundError:
    print(f"Could not find file: {FILENAME}")
except Exception as e:
    print(f"Error: {str(e)}")
