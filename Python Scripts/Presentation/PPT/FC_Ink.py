# Written by Ahmad Cooper
# Generates flashcards in PowerPoint from a word list. You can copy the word list in CSV format or separated by new lines.
import collections.abc # Essential
import comtypes.client
from pptx import Presentation
import pyperclip
import re
import time
import sys
import os
current_dir = os.path.dirname(os.path.abspath(__file__))
parent_dir = os.path.dirname(current_dir)
os.chdir(current_dir)
pdf_dir = os.path.join(parent_dir, "PDF")  
sys.path.append(f"{parent_dir}/PDF")
from RandomlyAddPDF2PDF import *        
from shuffle_2_pdf import *


clipboard = pyperclip.paste()
words = re.split(',|\n', clipboard)
final = [word.strip() for word in words if word.strip()]

# Define the name of the PowerPoint file
FILENAME = "PyFC"
ext = ".pptx"
try:
    # Try to open an existing PowerPoint presentation or create a new one
    prs = Presentation(f"Template/{FILENAME}{ext}")

    # Get the custom layout by ID
    custom_layout = prs.slide_layouts.get_by_name('3')
    if custom_layout is None:
        raise ValueError('Custom layout not found')
    

    for i in range(len(final)):
        # Add a new slide with the custom layout
        slide = prs.slides.add_slide(custom_layout)

        # Find the text box placeholder and set its text
        for shape in slide.placeholders:
            if shape.name == 'Text Placeholder 1':
                shape.text = str(final[i])
                break

        # Set the slide name
        slide.name = str(final[i])

    # Save the presentation
    prs.save(f"Output/{FILENAME}_1{ext}")  
    pdf = f"Output/{FILENAME}_1.pdf"

    powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
    powerpoint.Presentations.Open(os.path.abspath(f"Output/{FILENAME}_1{ext}"))
    powerpoint.ActivePresentation.SaveAs(os.path.abspath(pdf), 
                                      32) # 32 is the PDF format code
    powerpoint.Quit()
    print(f"Successfully created {len(final)} FC slides")
    shuffle_pdf(pdf)
    time.sleep(1)
    pdf = f"Output/{FILENAME}_1_shuffled.pdf"
    add_pdfs_to_pdf(pdf)
    os.startfile(f"{script_dir}/Output")

except FileNotFoundError:
    print(f"Could not find file: {FILENAME}")
except Exception as e:
    print(f"Error: {str(e)}")
